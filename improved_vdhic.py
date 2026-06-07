#!/usr/bin/env python3
"""
Create an improved VDHIC PowerPoint presentation using python-pptx
with enhanced design and professional styling
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.xmlchemy import OxmlElement

def create_improved_vdhic_presentation():
    # Create presentation
    prs = Presentation()

    # Define color scheme
    primary_blue = RGBColor(30, 58, 138)      # #1e3a8a
    light_blue = RGBColor(59, 130, 246)       # #3b82f6
    text_dark = RGBColor(31, 41, 55)          # #1f293b
    text_light = RGBColor(148, 163, 184)      # #94a3b8
    white = RGBColor(255, 255, 255)

    # Helper function to add gradient background
    def add_gradient_background(slide, start_color, end_color):
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135
        fill.gradient_stops[0].color.rgb = start_color
        fill.gradient_stops[1].color.rgb = end_color

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_background(slide1, primary_blue, light_blue)

    # Main title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Value-Based Digital Health\nInnovation Canvas"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = white
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "VDHIC Instruction Guide"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = text_light
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide1.shapes.add_textbox(Inches(1), Inches(4.8), Inches(7), Inches(0.6))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "From Problem to Impact: Navigating Indonesia's Digital Health Ecosystem"
    tagline_p = tagline_frame.paragraphs[0]
    tagline_p.font.size = Pt(18)
    tagline_p.font.color.rgb = text_light
    tagline_p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(7), Inches(0.4))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Featuring the Quintuple Aim Framework | Australia Awards Fellowship 2025 • Monash University"
    footer_p = footer_frame.paragraphs[0]
    footer_p.font.size = Pt(14)
    footer_p.font.color.rgb = text_light
    footer_p.alignment = PP_ALIGN.CENTER

    # Slide 2: What is VDHIC?
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Left panel
    left_panel = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4), Inches(7.5))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = primary_blue

    # Right panel background
    right_bg = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(0), Inches(6), Inches(7.5))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(248, 250, 252)

    # Left content
    left_title = slide2.shapes.add_textbox(Inches(0.3), Inches(1), Inches(3.4), Inches(1))
    left_title_frame = left_title.text_frame
    left_title_frame.text = "What is VDHIC?"
    left_title_p = left_title_frame.paragraphs[0]
    left_title_p.font.size = Pt(36)
    left_title_p.font.bold = True
    left_title_p.font.color.rgb = white

    left_desc = slide2.shapes.add_textbox(Inches(0.3), Inches(2.2), Inches(3.4), Inches(3))
    left_desc_frame = left_desc.text_frame
    left_desc_frame.text = "The Value-Based Digital Health Innovation Canvas transforms Indonesia's Digital Health Transformation Toolkit into a practical, actionable tool for innovators."
    for para in left_desc_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = white

    # Right content
    right_title = slide2.shapes.add_textbox(Inches(4.5), Inches(1), Inches(5), Inches(0.5))
    right_title_frame = right_title.text_frame
    right_title_frame.text = "Purpose"
    right_title_p = right_title_frame.paragraphs[0]
    right_title_p.font.size = Pt(24)
    right_title_p.font.bold = True
    right_title_p.font.color.rgb = primary_blue

    # Purpose items
    purpose_items = [
        "Structured framework for digital health innovation in Indonesia",
        "Embed governance considerations from the start",
        "Facilitate multi-stakeholder dialogue",
        "Bridge design to implementation gap",
        "Support regulatory sandbox applications",
        "Align with Quintuple Aim outcomes"
    ]

    for i, item in enumerate(purpose_items):
        item_box = slide2.shapes.add_textbox(Inches(4.5), Inches(1.8 + i*0.4), Inches(5), Inches(0.3))
        item_frame = item_box.text_frame
        item_frame.text = f"→ {item}"
        item_p = item_frame.paragraphs[0]
        item_p.font.size = Pt(14)
        item_p.font.color.rgb = text_dark

    # Slide 3: Canvas Structure
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = primary_blue

    header_title = slide3.shapes.add_textbox(Inches(0), Inches(0.2), Inches(10), Inches(0.5))
    header_title_frame = header_title.text_frame
    header_title_frame.text = "Canvas Structure"
    header_title_p = header_title_frame.paragraphs[0]
    header_title_p.font.size = Pt(36)
    header_title_p.font.bold = True
    header_title_p.font.color.rgb = white
    header_title_p.alignment = PP_ALIGN.CENTER

    header_subtitle = slide3.shapes.add_textbox(Inches(0), Inches(0.7), Inches(10), Inches(0.3))
    header_subtitle_frame = header_subtitle.text_frame
    header_subtitle_frame.text = "Four integrated layers forming a comprehensive innovation framework"
    header_subtitle_p = header_subtitle_frame.paragraphs[0]
    header_subtitle_p.font.size = Pt(16)
    header_subtitle_p.font.color.rgb = text_light
    header_subtitle_p.alignment = PP_ALIGN.CENTER

    # Grid items
    grid_data = [
        ("🛡️", "Governance Spine", "Four governance blocks forming the comprehensive compliance framework ensuring safety and regulatory adherence", light_blue, 0.5, 1.5),
        ("🎯", "Main Canvas Grid", "Nine blocks organized in three phases: Discover → Design → Deliver", RGBColor(16, 185, 129), 5.5, 1.5),
        ("⚖️", "Value Anchor", "Five Quintuple Aim outcome indicators measuring comprehensive value creation", RGBColor(245, 158, 11), 0.5, 3.5),
        ("🔄", "Innovation Phases", "Systematic progression from problem identification through solution design to market delivery", RGBColor(239, 68, 68), 5.5, 3.5)
    ]

    for icon, title, desc, color, x, y in grid_data:
        # Background box
        bg_box = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(4), Inches(1.3))
        bg_box.fill.solid()
        bg_box.fill.fore_color.rgb = white
        bg_box.line.color.rgb = color
        bg_box.line.width = Pt(2)

        # Accent bar
        accent = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.1), Inches(1.3))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color

        # Icon
        icon_box = slide3.shapes.add_textbox(Inches(x+0.1), Inches(y+0.1), Inches(0.5), Inches(0.4))
        icon_frame = icon_box.text_frame
        icon_frame.text = icon
        icon_p = icon_frame.paragraphs[0]
        icon_p.font.size = Pt(24)

        # Title
        title_box = slide3.shapes.add_textbox(Inches(x+0.7), Inches(y+0.1), Inches(3.2), Inches(0.3))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(16)
        title_p.font.bold = True
        title_p.font.color.rgb = primary_blue

        # Description
        desc_box = slide3.shapes.add_textbox(Inches(x+0.7), Inches(y+0.4), Inches(3.2), Inches(0.8))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_p = desc_frame.paragraphs[0]
        desc_p.font.size = Pt(11)
        desc_p.font.color.rgb = text_dark

    # Slide 4: Innovation Journey
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header4.fill.solid()
    header4.fill.fore_color.rgb = primary_blue

    header4_title = slide4.shapes.add_textbox(Inches(0), Inches(0.2), Inches(10), Inches(0.5))
    header4_title_frame = header4_title.text_frame
    header4_title_frame.text = "Innovation Journey"
    header4_title_p = header4_title_frame.paragraphs[0]
    header4_title_p.font.size = Pt(36)
    header4_title_p.font.bold = True
    header4_title_p.font.color.rgb = white
    header4_title_p.alignment = PP_ALIGN.CENTER

    # Phases
    phases_data = [
        ("🔍", "DISCOVER", ["Health Challenge", "User Profiles", "Stakeholder Map"], "Understanding the problem space before solutions", light_blue, 0.8, 1.5),
        ("🎨", "DESIGN", ["Solution Architecture", "User Experience", "SATUSEHAT Integration"], "Translating insights into concrete solutions", RGBColor(16, 185, 129), 3.5, 1.5),
        ("🚀", "DELIVER", ["Evidence Plan", "Scale Strategy", "Business Model"], "Validation, scaling, and sustainability", RGBColor(245, 158, 11), 6.2, 1.5)
    ]

    for icon, title, items, desc, color, x, y in phases_data:
        # Phase box
        phase_box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.4), Inches(4))
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = white
        phase_box.line.color.rgb = color
        phase_box.line.width = Pt(3)

        # Icon
        icon_box = slide4.shapes.add_textbox(Inches(x+0.9), Inches(y+0.2), Inches(0.6), Inches(0.4))
        icon_frame = icon_box.text_frame
        icon_frame.text = icon
        icon_p = icon_frame.paragraphs[0]
        icon_p.font.size = Pt(32)
        icon_p.alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide4.shapes.add_textbox(Inches(x), Inches(y+0.7), Inches(2.4), Inches(0.3))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(18)
        title_p.font.bold = True
        title_p.font.color.rgb = primary_blue
        title_p.alignment = PP_ALIGN.CENTER

        # Items
        for i, item in enumerate(items):
            item_box = slide4.shapes.add_textbox(Inches(x+0.2), Inches(y+1.2 + i*0.25), Inches(2), Inches(0.3))
            item_frame = item_box.text_frame
            item_frame.text = f"• {item}"
            item_p = item_frame.paragraphs[0]
            item_p.font.size = Pt(12)
            item_p.font.color.rgb = text_dark

        # Description box
        desc_box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y+3.3), Inches(2.4), Inches(0.7))
        desc_box.fill.solid()
        desc_box.fill.fore_color.rgb = RGBColor(248, 250, 252)

        desc_text_box = slide4.shapes.add_textbox(Inches(x+0.1), Inches(y+3.4), Inches(2.2), Inches(0.5))
        desc_text_frame = desc_text_box.text_frame
        desc_text_frame.text = desc
        desc_text_p = desc_text_frame.paragraphs[0]
        desc_text_p.font.size = Pt(10)
        desc_text_p.font.italic = True
        desc_text_p.font.color.rgb = text_light
        desc_text_p.alignment = PP_ALIGN.CENTER

    # Slide 5: Quintuple Aim
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header5.fill.solid()
    header5.fill.fore_color.rgb = primary_blue

    header5_title = slide5.shapes.add_textbox(Inches(0), Inches(0.2), Inches(10), Inches(0.5))
    header5_title_frame = header5_title.text_frame
    header5_title_frame.text = "The Quintuple Aim"
    header5_title_p = header5_title_frame.paragraphs[0]
    header5_title_p.font.size = Pt(36)
    header5_title_p.font.bold = True
    header5_title_p.font.color.rgb = white
    header5_title_p.alignment = PP_ALIGN.CENTER

    header5_subtitle = slide5.shapes.add_textbox(Inches(0), Inches(0.7), Inches(10), Inches(0.3))
    header5_subtitle_frame = header5_subtitle.text_frame
    header5_subtitle_frame.text = "Every innovation must demonstrate clear alignment with these five outcomes"
    header5_subtitle_p = header5_subtitle_frame.paragraphs[0]
    header5_subtitle_p.font.size = Pt(16)
    header5_subtitle_p.font.color.rgb = text_light
    header5_subtitle_p.alignment = PP_ALIGN.CENTER

    # Quintuple items
    quintuple_data = [
        ("🏥", "Population Health", "Clinical outcomes at population level", RGBColor(239, 68, 68), 0.6, 1.5),
        ("😊", "Patient Experience", "Satisfaction & engagement", light_blue, 2.04, 1.5),
        ("👨\u200d⚕️", "Provider Satisfaction", "Clinician workflow improvement", RGBColor(16, 185, 129), 3.48, 1.5),
        ("💰", "Cost Reduction", "Efficiency gains", RGBColor(245, 158, 11), 4.92, 1.5),
        ("⚖️", "Health Equity", "Reduce healthcare disparities", RGBColor(139, 92, 246), 6.36, 1.5)
    ]

    for icon, title, desc, color, x, y in quintuple_data:
        # Background
        bg_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(1.2), Inches(3))
        bg_box.fill.solid()
        bg_box.fill.fore_color.rgb = white
        bg_box.line.color.rgb = color
        bg_box.line.width = Pt(2)

        # Icon
        icon_box = slide5.shapes.add_textbox(Inches(x+0.35), Inches(y+0.2), Inches(0.5), Inches(0.4))
        icon_frame = icon_box.text_frame
        icon_frame.text = icon
        icon_p = icon_frame.paragraphs[0]
        icon_p.font.size = Pt(28)
        icon_p.alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide5.shapes.add_textbox(Inches(x), Inches(y+0.7), Inches(1.2), Inches(0.4))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(12)
        title_p.font.bold = True
        title_p.font.color.rgb = primary_blue
        title_p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide5.shapes.add_textbox(Inches(x+0.05), Inches(y+1.2), Inches(1.1), Inches(0.8))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_p = desc_frame.paragraphs[0]
        desc_p.font.size = Pt(10)
        desc_p.font.color.rgb = text_dark
        desc_p.alignment = PP_ALIGN.CENTER
        desc_p.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Save presentation
    output_file = 'VDHIC_Improved_Presentation.pptx'
    prs.save(output_file)
    print(f"Improved presentation created successfully: {output_file}")

if __name__ == "__main__":
    create_improved_vdhic_presentation()