#!/usr/bin/env python3
"""
Convert VDHIC HTML presentation to PowerPoint format
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_vdhic_presentation():
    # Create presentation
    prs = Presentation()

    # Define color scheme
    primary_color = RGBColor(102, 126, 234)  # #667eea
    secondary_color = RGBColor(118, 75, 162)  # #764ba2
    text_color = RGBColor(51, 51, 51)  # #333333
    white = RGBColor(255, 255, 255)

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box1 = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    title_frame1 = title_box1.text_frame
    title_frame1.text = "Value-Based Digital Health"
    title_p1 = title_frame1.paragraphs[0]
    title_p1.font.size = Pt(48)
    title_p1.font.bold = True
    title_p1.font.color.rgb = primary_color
    title_p1.alignment = PP_ALIGN.CENTER

    # Title 2
    title_box2 = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "Innovation Canvas"
    title_p2 = title_frame2.paragraphs[0]
    title_p2.font.size = Pt(48)
    title_p2.font.bold = True
    title_p2.font.color.rgb = primary_color
    title_p2.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "VDHIC - Instruction Guide"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = secondary_color
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Content box
    content_box = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(6), Inches(7), Inches(2.5))
    content_fill = content_box.fill
    content_fill.solid()
    content_fill.fore_color.rgb = RGBColor(245, 245, 245)
    content_line = content_box.line
    content_line.color.rgb = primary_color
    content_line.width = Pt(2)

    content_text = content_box.text_frame
    content_text.text = "From Problem to Impact: Navigating Indonesia's Digital Health Ecosystem\n\nFeaturing the Quintuple Aim Framework\n\nAustralia Awards Fellowship 2025 • Monash University"
    content_text.margin_left = Inches(0.2)
    content_text.margin_right = Inches(0.2)
    content_text.margin_top = Inches(0.2)
    content_text.margin_bottom = Inches(0.2)

    for paragraph in content_text.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.alignment = PP_ALIGN.CENTER
        if paragraph.text.startswith("From"):
            paragraph.font.bold = True

    # Slide 2: Introduction
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title2_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title2_frame = title2_box.text_frame
    title2_frame.text = "What is VDHIC?"
    title2_p = title2_frame.paragraphs[0]
    title2_p.font.size = Pt(44)
    title2_p.font.bold = True
    title2_p.font.color.rgb = primary_color
    title2_p.alignment = PP_ALIGN.CENTER

    # Main content
    main_box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(8), Inches(7))
    main_box.fill.solid()
    main_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    main_box.line.color.rgb = primary_color
    main_box.line.width = Pt(1)

    main_text = main_box.text_frame
    main_text.text = "The Value-Based Digital Health Innovation Canvas (VDHIC) transforms the Indonesia Digital Health Transformation Toolkit into a practical, actionable tool for innovators.\n\nPurpose:\n• Structured framework for digital health innovation in Indonesia\n• Embed governance considerations from the start\n• Facilitate multi-stakeholder dialogue\n• Bridge design to implementation gap\n• Support regulatory sandbox applications\n• Align with Quintuple Aim outcomes"

    for paragraph in main_text.paragraphs:
        paragraph.font.size = Pt(20)
        if "Purpose:" in paragraph.text:
            paragraph.font.bold = True
            paragraph.font.color.rgb = primary_color
        elif paragraph.text.startswith("•"):
            paragraph.font.size = Pt(18)
            paragraph.level = 1

    # Slide 3: Canvas Structure
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title3_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title3_frame = title3_box.text_frame
    title3_frame.text = "Canvas Structure"
    title3_p = title3_frame.paragraphs[0]
    title3_p.font.size = Pt(44)
    title3_p.font.bold = True
    title3_p.font.color.rgb = primary_color
    title3_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle3_box = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    subtitle3_frame = subtitle3_box.text_frame
    subtitle3_frame.text = "Four integrated layers forming a comprehensive innovation framework"
    subtitle3_p = subtitle3_frame.paragraphs[0]
    subtitle3_p.font.size = Pt(20)
    subtitle3_p.alignment = PP_ALIGN.CENTER

    # Grid items
    items = [
        ("🛡️ Governance Spine", "Four governance blocks forming the comprehensive compliance framework ensuring safety and regulatory adherence"),
        ("🎯 Main Canvas Grid", "Nine blocks organized in three phases: Discover → Design → Deliver"),
        ("⚖️ Value Anchor", "Five Quintuple Aim outcome indicators measuring comprehensive value creation"),
        ("🔄 Innovation Phases", "Systematic progression from problem identification through solution design to market delivery")
    ]

    positions = [(1, 2.5), (5.5, 2.5), (1, 5), (5.5, 5)]

    for i, (title, content) in enumerate(items):
        left, top = positions[i]
        grid_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(4), Inches(2))

        # Gradient effect simulation
        grid_box.fill.solid()
        if i % 2 == 0:
            grid_box.fill.fore_color.rgb = primary_color
        else:
            grid_box.fill.fore_color.rgb = secondary_color

        grid_text = grid_box.text_frame
        grid_text.text = f"{title}\n\n{content}"
        grid_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        for paragraph in grid_text.paragraphs:
            if title in paragraph.text:
                paragraph.font.size = Pt(18)
                paragraph.font.bold = True
                paragraph.font.color.rgb = white
            else:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = white
            paragraph.alignment = PP_ALIGN.CENTER

    # Slide 4: Three Phases
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title4_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title4_frame = title4_box.text_frame
    title4_frame.text = "Innovation Journey"
    title4_p = title4_frame.paragraphs[0]
    title4_p.font.size = Pt(44)
    title4_p.font.bold = True
    title4_p.font.color.rgb = primary_color
    title4_p.alignment = PP_ALIGN.CENTER

    # Phases
    phases = [
        ("🔍 DISCOVER", ["Health Challenge", "User Profiles", "Stakeholder Map"], "Understanding the problem space before solutions"),
        ("🎨 DESIGN", ["Solution Architecture", "User Experience", "SATUSEHAT Integration"], "Translating insights into concrete solutions"),
        ("🚀 DELIVER", ["Evidence Plan", "Scale Strategy", "Business Model"], "Validation, scaling, and sustainability")
    ]

    phase_positions = [(0.5, 2), (3.5, 2), (6.5, 2)]

    for i, (title, items, description) in enumerate(phases):
        left, top = phase_positions[i]

        # Phase box
        phase_box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(3), Inches(4))
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = white
        phase_box.line.color.rgb = primary_color
        phase_box.line.width = Pt(2)

        # Phase title
        phase_title_box = slide4.shapes.add_textbox(Inches(left+0.1), Inches(top+0.2), Inches(2.8), Inches(0.5))
        phase_title_frame = phase_title_box.text_frame
        phase_title_frame.text = title
        phase_title_p = phase_title_frame.paragraphs[0]
        phase_title_p.font.size = Pt(20)
        phase_title_p.font.bold = True
        phase_title_p.font.color.rgb = primary_color
        phase_title_p.alignment = PP_ALIGN.CENTER

        # Phase items
        items_box = slide4.shapes.add_textbox(Inches(left+0.2), Inches(top+0.8), Inches(2.6), Inches(2))
        items_frame = items_box.text_frame
        items_text = ""
        for item in items:
            items_text += f"• {item}\n"
        items_frame.text = items_text.rstrip()

        for paragraph in items_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.level = 0

        # Description
        desc_box = slide4.shapes.add_textbox(Inches(left+0.1), Inches(top+3), Inches(2.8), Inches(0.8))
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_p = desc_frame.paragraphs[0]
        desc_p.font.size = Pt(14)
        desc_p.font.color.rgb = RGBColor(102, 102, 102)
        desc_p.alignment = PP_ALIGN.CENTER
        desc_p.font.italic = True

    # Slide 5: Quintuple Aim
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title5_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title5_frame = title5_box.text_frame
    title5_frame.text = "The Quintuple Aim"
    title5_p = title5_frame.paragraphs[0]
    title5_p.font.size = Pt(44)
    title5_p.font.bold = True
    title5_p.font.color.rgb = primary_color
    title5_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle5_box = slide5.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    subtitle5_frame = subtitle5_box.text_frame
    subtitle5_frame.text = "Every innovation must demonstrate clear alignment with these five outcomes"
    subtitle5_p = subtitle5_frame.paragraphs[0]
    subtitle5_p.font.size = Pt(20)
    subtitle5_p.alignment = PP_ALIGN.CENTER

    # Quintuple items
    quintuple = [
        ("🏥 Population\nHealth", "Clinical outcomes at population level", (245, 87, 108)),
        ("😊 Patient\nExperience", "Satisfaction & engagement", (79, 172, 254)),
        ("👨‍⚕️ Provider\nSatisfaction", "Clinician workflow improvement", (67, 233, 123)),
        ("💰 Cost\nReduction", "Efficiency gains", (250, 112, 154)),
        ("⚖️ Health\nEquity", "Reduce healthcare disparities", (48, 207, 208))
    ]

    quintuple_positions = [(0.5, 2.5), (2.1, 2.5), (3.7, 2.5), (5.3, 2.5), (6.9, 2.5)]

    for i, (title, content, color) in enumerate(quintuple):
        left, top = quintuple_positions[i]

        # Quintuple box
        quintuple_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(1.5), Inches(3))
        quintuple_box.fill.solid()
        quintuple_box.fill.fore_color.rgb = RGBColor(*color)

        quintuple_text = quintuple_box.text_frame
        quintuple_text.text = f"{title}\n\n{content}"
        quintuple_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        for paragraph in quintuple_text.paragraphs:
            if "Health" in paragraph.text or "Experience" in paragraph.text or "Satisfaction" in paragraph.text or "Reduction" in paragraph.text or "Equity" in paragraph.text:
                paragraph.font.size = Pt(16)
                paragraph.font.bold = True
                paragraph.font.color.rgb = white
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = white
                paragraph.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save('VDHIC_Instruction_Guide_Presentation.pptx')
    print("Presentation created successfully: VDHIC_Instruction_Guide_Presentation.pptx")

if __name__ == "__main__":
    create_vdhic_presentation()